#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build a Workshop Design System PowerPoint TEMPLATE (.pptx).

A 15-slide, bilingual (English + Hebrew RTL) reference deck that carries every
core building block of the design system as *native, editable* PowerPoint
objects — the six-ink color system, the type scale, text/label styles, an icon
library (rendered from the real SVG sprite), KPI cards, problem/solution,
process flow, native charts, a comparison table, roadmap, risks, a 2x2 matrix,
and a closing slide.

Intended to be reused as a starting template, or handed to another AI agent so
it can generate on-brand slides using the correct colors, fonts and metrics.

Usage:  python3 scripts/build-pptx-template.py [out.pptx] [--icons DIR]

Note: PowerPoint cannot reproduce the hand-drawn SVG stroke or the sketchbook
hand-print font (those are HTML/CSS/JS only). This template therefore uses the
system's *clean* brand mode: Fraunces/Frank Ruhl Libre display, Inter/Heebo
body, IBM Plex Mono, thin ink outlines standing in for the sketch border.
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ---------------------------------------------------------------- tokens ---
# Six-ink semantic palette (light mode) — from project/css/tokens.css
INK = {
    "black":  RGBColor(0x1B, 0x1E, 0x23),
    "blue":   RGBColor(0x1E, 0x5A, 0x8C),
    "red":    RGBColor(0xB0, 0x33, 0x2E),
    "green":  RGBColor(0x3F, 0x7A, 0x44),
    "orange": RGBColor(0xC1, 0x63, 0x1A),
    "purple": RGBColor(0x6B, 0x4E, 0x9E),
}
INK_MEANING = {
    "black":  ("Structure",        "מבנה"),
    "blue":   ("Future / ideas",   "עתיד / רעיונות"),
    "red":    ("Problems / pain",  "בעיות / מצב קיים"),
    "green":  ("Solutions / done", "פתרונות / הושלם"),
    "orange": ("Risks / progress", "סיכונים / בתהליך"),
    "purple": ("AI / automation",  "AI / אוטומציה"),
}
INK_HEX = {"black":"#1B1E23","blue":"#1E5A8C","red":"#B0332E",
           "green":"#3F7A44","orange":"#C1631A","purple":"#6B4E9E"}

PAPER        = RGBColor(0xFB, 0xFA, 0xF7)   # --color-paper
PAPER_RAISED = RGBColor(0xFF, 0xFF, 0xFF)   # --color-paper-raised
PAPER_SUNKEN = RGBColor(0xF3, 0xF1, 0xEC)   # --color-paper-sunken
TXT   = RGBColor(0x1B, 0x1E, 0x23)          # --color-ink
TXT2  = RGBColor(0x40, 0x45, 0x4E)          # --color-ink-2
TXT3  = RGBColor(0x6B, 0x70, 0x78)          # --color-ink-3
ACCENT= RGBColor(0x1B, 0x4B, 0x66)          # --color-accent

# Fonts (Google/self-hosted in the web system; install these for the .pptx)
F_DISP_LAT = "Fraunces"
F_DISP_HE  = "Frank Ruhl Libre"
F_BODY_LAT = "Inter"
F_BODY_HE  = "Heebo"
F_MONO     = "IBM Plex Mono"
F_HAND     = "Caveat"

# 16:9 canvas
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

# Icons are committed under assets/pptx-icons/ (regenerate with
# scripts/render-pptx-icons.mjs). Falls back to that repo path by default.
ICON_DIR = Path(__file__).resolve().parent.parent / "assets" / "pptx-icons"

# ---------------------------------------------------------------- helpers ---
def _set_cs(run, typeface):
    """Set the complex-script (Hebrew) + latin typeface on a run."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", typeface)

def _rtl(para):
    para._p.get_or_add_pPr().set("rtl", "1")

def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             rtl=False, line_spacing=None, space_after=None, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of (text, dict) runs.
       dict keys: size, color, font, bold, italic, tracking(char spacing pt)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if rtl:
            _rtl(p)
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = Pt(space_after)
        for text, st in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(st.get("size", 16))
            r.font.bold = st.get("bold", False)
            r.font.italic = st.get("italic", False)
            r.font.color.rgb = st.get("color", TXT)
            fnt = st.get("font", F_BODY_LAT)
            r.font.name = fnt
            _set_cs(r, fnt)
            if "tracking" in st:
                r._r.get_or_add_rPr().set("spc", str(int(st["tracking"] * 100)))
    return tb

def add_card(slide, x, y, w, h, ink="black", fill=PAPER_RAISED, line_w=1.75,
             radius=True):
    """Rounded rectangle with a thin ink outline — the clean-mode stand-in for
       the hand-drawn sketch border."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = INK[ink]
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.045
        except Exception:
            pass
    return shp

def add_rect(slide, x, y, w, h, fill, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def add_oval(slide, x, y, d, ink="black", fill=PAPER_RAISED, line_w=1.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = INK[ink]
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def add_icon(slide, name, x, y, size):
    p = ICON_DIR / (name + ".png")
    if p.exists():
        return slide.shapes.add_picture(str(p), x, y, size, size)
    return None

def bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def label(slide, x, y, w, text_lat, text_he=None, ink="black"):
    """Uppercase-tracked Latin label + optional plain Hebrew label."""
    runs = [[(text_lat.upper(), {"size": 11, "color": INK[ink], "bold": True,
                                 "font": F_BODY_LAT, "tracking": 1.5})]]
    add_text(slide, x, y, w, Inches(0.3), runs, align=PP_ALIGN.LEFT)
    if text_he:
        add_text(slide, x, y + Inches(0.24), w, Inches(0.3),
                 [[(text_he, {"size": 12, "color": INK[ink], "bold": True,
                              "font": F_BODY_HE})]],
                 align=PP_ALIGN.LEFT, rtl=True)

def slide_title(slide, en, he, kicker=None):
    """Standard slide header: English display title + Hebrew subtitle."""
    add_text(slide, Inches(0.6), Inches(0.42), Inches(12.1), Inches(0.9),
             [[(en, {"size": 34, "color": TXT, "bold": True, "font": F_DISP_LAT})]],
             align=PP_ALIGN.LEFT)
    add_text(slide, Inches(0.6), Inches(1.18), Inches(12.1), Inches(0.5),
             [[(he, {"size": 20, "color": ACCENT, "font": F_DISP_HE})]],
             align=PP_ALIGN.RIGHT, rtl=True)
    if kicker:
        add_text(slide, Inches(0.62), Inches(0.16), Inches(12), Inches(0.3),
                 [[(kicker.upper(), {"size": 10.5, "color": TXT3, "bold": True,
                                     "font": F_BODY_LAT, "tracking": 2})]])
    # thin rule under header
    add_rect(slide, Inches(0.6), Inches(1.62), Inches(12.13), Pt(1.5), TXT2)

def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg(s)
    return s

# ================================================================ build ===
def build(out_path):
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    IN = Inches

    # ---- 1 · COVER ----------------------------------------------------
    s = new_slide(prs)
    add_rect(s, 0, 0, EMU_W, IN(0.28), ACCENT)                    # top band
    add_icon(s, "icon-cpu", IN(0.7), IN(0.9), IN(1.0))
    add_text(s, IN(0.7), IN(2.4), IN(12), IN(1.6),
             [[("Workshop Design System", {"size": 58, "color": TXT, "bold": True,
                                           "font": F_DISP_LAT})]])
    add_text(s, IN(0.7), IN(3.7), IN(12), IN(0.8),
             [[("מערכת העיצוב — תבנית מצגת מנהלים", {"size": 30, "color": ACCENT,
                                                     "font": F_DISP_HE})]],
             align=PP_ALIGN.RIGHT, rtl=True)
    add_text(s, IN(0.72), IN(4.7), IN(12), IN(0.5),
             [[("PowerPoint template · English + עברית · six-ink system · "
                "editable objects",
                {"size": 15, "color": TXT3, "font": F_BODY_LAT})]])
    # six ink dots along the bottom
    dx = IN(0.72)
    for k in ["black", "blue", "red", "green", "orange", "purple"]:
        add_oval(s, dx, IN(6.4), IN(0.34), ink=k, fill=INK[k], line_w=0)
        dx += IN(0.55)
    add_text(s, IN(4.2), IN(6.42), IN(8.8), IN(0.4),
             [[("15-slide element gallery — duplicate any slide and edit in place",
                {"size": 12, "color": TXT3, "italic": True, "font": F_BODY_LAT})]],
             align=PP_ALIGN.RIGHT)

    # ---- 2 · HOW TO USE ----------------------------------------------
    s = new_slide(prs)
    slide_title(s, "How to use this template", "כיצד להשתמש בתבנית", kicker="Start here")
    colw = IN(5.9)
    en = [
        [("This deck is a component gallery. ", {"size": 15, "bold": True, "color": TXT}),
         ("Duplicate any slide and replace its content — the colors, fonts and "
          "layouts are already on-brand.", {"size": 15, "color": TXT2})],
        [("Six inks are semantic", {"size": 15, "bold": True, "color": TXT}),
         (" — never remap them: black=structure, blue=future, red=problems, "
          "green=solutions, orange=risks, purple=AI.", {"size": 15, "color": TXT2})],
        [("Fonts to install: ", {"size": 15, "bold": True, "color": TXT}),
         ("Fraunces, Frank Ruhl Libre, Inter, Heebo, IBM Plex Mono, Caveat "
          "(all free / Google Fonts).", {"size": 15, "color": TXT2})],
        [("No emoji, no gradients. ", {"size": 15, "bold": True, "color": TXT}),
         ("Numbers and Latin terms stay LTR inside Hebrew text.", {"size": 15, "color": TXT2})],
        [("Hand-sketch mode is HTML-only", {"size": 15, "bold": True, "color": INK["orange"]}),
         (" — this .pptx uses the clean brand mode instead.", {"size": 15, "color": TXT2})],
    ]
    he = [
        [("המצגת הזו היא גלריית רכיבים. ", {"size": 15, "bold": True, "color": TXT, "font": F_BODY_HE}),
         ("שכפלו כל שקף והחליפו תוכן — הצבעים, הפונטים והפריסות כבר במותג.",
          {"size": 15, "color": TXT2, "font": F_BODY_HE})],
        [("שש הדיו סמנטיות", {"size": 15, "bold": True, "color": TXT, "font": F_BODY_HE}),
         (" — אין למפות מחדש: שחור=מבנה, כחול=עתיד, אדום=בעיות, ירוק=פתרונות, "
          "כתום=סיכונים, סגול=AI.", {"size": 15, "color": TXT2, "font": F_BODY_HE})],
        [("פונטים להתקנה: ", {"size": 15, "bold": True, "color": TXT, "font": F_BODY_HE}),
         ("Fraunces, Frank Ruhl Libre, Inter, Heebo, IBM Plex Mono, Caveat.",
          {"size": 15, "color": TXT2, "font": F_BODY_HE})],
        [("בלי אימוג'י, בלי גרדיאנטים. ", {"size": 15, "bold": True, "color": TXT, "font": F_BODY_HE}),
         ("מספרים ומונחים לועזיים נשארים LTR בתוך טקסט עברי.",
          {"size": 15, "color": TXT2, "font": F_BODY_HE})],
        [("מצב הסקיצה קיים רק ב-HTML", {"size": 15, "bold": True, "color": INK["orange"], "font": F_BODY_HE}),
         (" — כאן משתמשים במצב המותג הנקי.", {"size": 15, "color": TXT2, "font": F_BODY_HE})],
    ]
    c1 = add_card(s, IN(0.6), IN(1.9), colw, IN(5.0), ink="blue")
    label(s, IN(0.9), IN(2.1), IN(5), "English")
    add_text(s, IN(0.9), IN(2.55), IN(5.3), IN(4.2), en, space_after=8, line_spacing=1.05)
    c2 = add_card(s, IN(6.83), IN(1.9), colw, IN(5.0), ink="black")
    label(s, IN(11.5), IN(2.1), IN(1.0), "עברית")
    add_text(s, IN(7.1), IN(2.55), IN(5.3), IN(4.2), he, align=PP_ALIGN.RIGHT,
             rtl=True, space_after=8, line_spacing=1.05)

    # ---- 3 · COLOR SYSTEM --------------------------------------------
    s = new_slide(prs)
    slide_title(s, "The six-ink color system", "מערכת שש הדיו", kicker="Foundations")
    order = ["black", "blue", "red", "green", "orange", "purple"]
    cw, ch = IN(3.9), IN(2.35)
    for i, k in enumerate(order):
        col, row = i % 3, i // 3
        x = IN(0.6) + col * IN(4.13)
        y = IN(1.95) + row * IN(2.55)
        add_card(s, x, y, cw, ch, ink=k)
        add_rect(s, x + IN(0.18), y + IN(0.18), IN(1.0), IN(1.0), INK[k])  # swatch
        en_n, he_n = INK_MEANING[k]
        add_text(s, x + IN(1.35), y + IN(0.16), IN(2.4), IN(0.5),
                 [[(k.capitalize(), {"size": 19, "bold": True, "color": INK[k],
                                     "font": F_DISP_LAT})]])
        add_text(s, x + IN(1.35), y + IN(0.62), IN(2.4), IN(0.4),
                 [[(INK_HEX[k], {"size": 13, "color": TXT3, "font": F_MONO})]])
        add_text(s, x + IN(0.2), y + IN(1.35), IN(3.5), IN(0.4),
                 [[(en_n, {"size": 14, "bold": True, "color": TXT})]])
        add_text(s, x + IN(0.2), y + IN(1.78), IN(3.5), IN(0.4),
                 [[(he_n, {"size": 15, "color": TXT2, "font": F_BODY_HE})]],
                 align=PP_ALIGN.RIGHT, rtl=True)

    # ---- 4 · TYPOGRAPHY ----------------------------------------------
    s = new_slide(prs)
    slide_title(s, "Typography", "טיפוגרפיה", kicker="Foundations")
    rows = [
        ("Display", "Fraunces / Frank Ruhl Libre", 40,
         "Executive headline", "כותרת ניהולית", F_DISP_LAT, F_DISP_HE, TXT),
        ("Body", "Inter / Heebo", 20,
         "Readable running text for content.", "טקסט גוף קריא לתוכן השקף.",
         F_BODY_LAT, F_BODY_HE, TXT2),
        ("Label", "Inter · uppercase · tracked", 13,
         "CURRENT SITUATION", "מצב נוכחי", F_BODY_LAT, F_BODY_HE, INK["red"]),
        ("Mono", "IBM Plex Mono · figures", 20,
         "1,247  +18%  -1.3", "‎12,400", F_MONO, F_MONO, INK["blue"]),
        ("Hand", "Caveat (Latin annotations only)", 26,
         "a quick note in the margin", "— (hand font: Latin only)",
         F_HAND, F_BODY_HE, INK["purple"]),
    ]
    y = IN(1.95)
    for name, meta, size, en_s, he_s, fl, fh, col in rows:
        add_text(s, IN(0.6), y + IN(0.05), IN(2.2), IN(0.5),
                 [[(name, {"size": 16, "bold": True, "color": TXT, "font": F_BODY_LAT})]])
        add_text(s, IN(0.6), y + IN(0.5), IN(2.6), IN(0.4),
                 [[(meta, {"size": 10.5, "color": TXT3, "font": F_MONO})]])
        add_text(s, IN(3.0), y, IN(5.4), IN(1.0),
                 [[(en_s, {"size": size, "color": col, "font": fl,
                           "bold": name in ("Display",)})]],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, IN(8.5), y, IN(4.2), IN(1.0),
                 [[(he_s, {"size": size, "color": col, "font": fh,
                           "bold": name in ("Display",)})]],
                 align=PP_ALIGN.RIGHT, rtl=True, anchor=MSO_ANCHOR.MIDDLE)
        y += IN(1.02)
        add_rect(s, IN(0.6), y - IN(0.06), IN(12.13), Pt(0.75), PAPER_SUNKEN)

    # ---- 5 · TEXT STYLES & LABELS ------------------------------------
    s = new_slide(prs)
    slide_title(s, "Text styles, labels & numbers", "סגנונות טקסט, תוויות ומספרים",
                kicker="Foundations")
    # left: label styles
    add_card(s, IN(0.6), IN(1.95), IN(5.9), IN(5.0), ink="black")
    label(s, IN(0.85), IN(2.15), IN(5), "Label styles", "סגנונות תווית")
    ex = [
        [("Uppercase tracked (Latin): ", {"size": 13, "color": TXT2}),
         ("CURRENT SITUATION", {"size": 13, "bold": True, "color": TXT, "tracking": 1.5})],
        [("Hebrew plain label: ", {"size": 13, "color": TXT2}),
         ("מצב נוכחי", {"size": 14, "bold": True, "color": TXT, "font": F_BODY_HE})],
        [("Section kicker: ", {"size": 13, "color": TXT2}),
         ("FOUNDATIONS · 01", {"size": 12, "bold": True, "color": TXT3, "tracking": 2})],
        [("Emphasis — highlight: ", {"size": 13, "color": TXT2}),
         ("key takeaway", {"size": 13, "bold": True, "color": TXT,
          "font": F_HAND, "size": 16})],
        [("Emphasis — circle a number: ", {"size": 13, "color": TXT2}),
         ("78", {"size": 15, "bold": True, "color": INK["red"]})],
        [("Rejected idea (strike): ", {"size": 13, "color": TXT2}),
         ("old process", {"size": 13, "italic": True, "color": TXT3})],
    ]
    add_text(s, IN(0.85), IN(2.75), IN(5.4), IN(4), ex, space_after=12, line_spacing=1.05)
    # right: numbers & deltas
    add_card(s, IN(6.83), IN(1.95), IN(5.9), IN(5.0), ink="blue")
    label(s, IN(7.08), IN(2.15), IN(5), "Numbers & deltas", "מספרים ודלתאות")
    nums = [
        ("Big KPI value", "1,247", INK["black"], "24"),
        ("Positive delta (good)", "+18%", INK["green"], "18"),
        ("Negative delta (bad)", "-9%", INK["red"], "18"),
        ("Negative-is-good", "-1.3 days", INK["green"], "18"),
        ("Hebrew + LTR number", "זמן מענה 9 ימים", INK["black"], "16"),
    ]
    yy = IN(2.8)
    for cap, val, col, sz in nums:
        add_text(s, IN(7.08), yy, IN(3.2), IN(0.5),
                 [[(cap, {"size": 12, "color": TXT2})]], anchor=MSO_ANCHOR.MIDDLE)
        is_he = any('֐' <= c <= '׿' for c in val)
        add_text(s, IN(10.2), yy, IN(2.3), IN(0.5),
                 [[(val, {"size": int(sz), "bold": True, "color": col,
                          "font": F_BODY_HE if is_he else F_MONO})]],
                 align=PP_ALIGN.RIGHT, rtl=is_he, anchor=MSO_ANCHOR.MIDDLE)
        yy += IN(0.72)
    add_text(s, IN(7.08), IN(6.35), IN(5.4), IN(0.5),
             [[("Rule: color deltas by meaning, not by sign.",
                {"size": 12, "italic": True, "color": TXT3})]])

    # ---- 6 · ICON LIBRARY --------------------------------------------
    s = new_slide(prs)
    slide_title(s, "Icon library", "ספריית אייקונים", kicker="Assets · 70+ glyphs")
    grid = ["icon-mail","icon-phone","icon-clock","icon-user","icon-users","icon-chart-bar",
            "icon-chart-pie","icon-target","icon-trend-up","icon-trend-down","icon-cpu","icon-robot",
            "icon-sparkle","icon-shield-check","icon-warning","icon-check","icon-rocket","icon-lightbulb",
            "icon-calendar","icon-gear","icon-database","icon-globe","icon-search","icon-flag",
            "icon-star","icon-handshake","icon-building","icon-map-pin","icon-file-text","icon-lock",
            "icon-key","icon-trophy","icon-route","icon-layers","icon-gauge","icon-list"]
    cols = 9
    cellw, cellh = IN(1.36), IN(1.28)
    x0, y0 = IN(0.62), IN(2.0)
    for i, name in enumerate(grid):
        col, row = i % cols, i // cols
        x = x0 + col * cellw
        y = y0 + row * cellh
        add_icon(s, name, x + IN(0.34), y, IN(0.68))
        add_text(s, x - IN(0.05), y + IN(0.72), cellw, IN(0.4),
                 [[(name.replace("icon-", ""), {"size": 8.5, "color": TXT3,
                                                "font": F_MONO})]],
                 align=PP_ALIGN.CENTER)
    add_text(s, IN(0.62), IN(6.9), IN(12), IN(0.4),
             [[("Own stroke set · 24×24 · round caps · currentColor · hatch-fill "
                "duotone variants · reference: <use href=\"#icon-name\">",
                {"size": 11, "color": TXT3, "italic": True})]])

    # ---- 7 · KPI CARDS -----------------------------------------------
    s = new_slide(prs)
    slide_title(s, "KPI cards", "כרטיסי מדד", kicker="Components")
    kpis = [
        ("icon-mail",  "red",    "פניות בחודש",        "12,400", "+18%",  "vs last year",  "מול שנה שעברה"),
        ("icon-clock", "red",    "זמן מענה ממוצע",     "9 ימים", "+1.5",  "vs last Q",     "מול רבעון קודם"),
        ("icon-phone", "orange", "שיחות ננטשות",       "31%",    "+6%",   "vs last Q",     "מול רבעון קודם"),
        ("icon-trend-up","green","שביעות רצון",        "80%",    "+18",   "target",        "יעד"),
    ]
    cw = IN(2.95)
    for i, (ic, ink, he_lab, val, delta, en_f, he_f) in enumerate(kpis):
        x = IN(0.6) + i * IN(3.06)
        add_card(s, x, IN(2.4), cw, IN(3.4), ink=ink)
        add_icon(s, ic, x + IN(0.25), IN(2.62), IN(0.6))
        add_text(s, x + IN(0.2), IN(3.35), cw - IN(0.4), IN(0.5),
                 [[(he_lab, {"size": 15, "bold": True, "color": INK[ink],
                             "font": F_BODY_HE})]], align=PP_ALIGN.RIGHT, rtl=True)
        is_he = any('֐' <= c <= '׿' for c in val)
        add_text(s, x + IN(0.2), IN(3.85), cw - IN(0.4), IN(0.9),
                 [[(val, {"size": 40, "bold": True, "color": TXT,
                          "font": F_BODY_HE if is_he else F_MONO})]],
                 align=PP_ALIGN.RIGHT if is_he else PP_ALIGN.LEFT, rtl=is_he)
        add_text(s, x + IN(0.2), IN(4.75), cw - IN(0.4), IN(0.4),
                 [[(delta, {"size": 16, "bold": True, "color": INK[ink],
                            "font": F_MONO})]])
        add_text(s, x + IN(0.2), IN(5.25), cw - IN(0.4), IN(0.4),
                 [[(he_f, {"size": 11, "color": TXT3, "font": F_BODY_HE})]],
                 align=PP_ALIGN.RIGHT, rtl=True)

    # ---- 8 · PROBLEM ↔ SOLUTION --------------------------------------
    s = new_slide(prs)
    slide_title(s, "Problem ↔ Solution", "בעיה מול פתרון", kicker="Components")
    add_card(s, IN(0.6), IN(2.1), IN(5.5), IN(4.6), ink="red")
    label(s, IN(0.9), IN(2.3), IN(5), "Current — pain", "היום — הכאב")
    prob = [[("• תהליך ידני", {"size": 17, "color": TXT, "font": F_BODY_HE})],
            [("• חוסר נראות", {"size": 17, "color": TXT, "font": F_BODY_HE})],
            [("• פניות הולכות לאיבוד", {"size": 17, "color": TXT, "font": F_BODY_HE})],
            [("• אין דוחות להנהלה", {"size": 17, "color": TXT, "font": F_BODY_HE})]]
    add_text(s, IN(0.9), IN(3.1), IN(5), IN(3.2), prob, align=PP_ALIGN.RIGHT,
             rtl=True, space_after=10)
    # arrow
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, IN(6.25), IN(4.1), IN(0.85), IN(0.6))
    ar.rotation = 180
    ar.fill.solid(); ar.fill.fore_color.rgb = INK["green"]
    ar.line.fill.background(); ar.shadow.inherit = False
    add_card(s, IN(7.23), IN(2.1), IN(5.5), IN(4.6), ink="blue")
    label(s, IN(7.53), IN(2.3), IN(5), "Future — solution", "מחר — הפתרון")
    sol = [[("✓ פלטפורמה מרכזית אחת", {"size": 17, "color": TXT, "font": F_BODY_HE})],
           [("✓ סיווג בסיוע AI", {"size": 17, "color": TXT, "font": F_BODY_HE})],
           [("✓ ניתוב אוטומטי לבעלים", {"size": 17, "color": TXT, "font": F_BODY_HE})],
           [("✓ דשבורד ב-Real-time", {"size": 17, "color": TXT, "font": F_BODY_HE})]]
    add_text(s, IN(7.53), IN(3.1), IN(5), IN(3.2), sol, align=PP_ALIGN.RIGHT,
             rtl=True, space_after=10)

    # ---- 9 · PROCESS FLOW --------------------------------------------
    s = new_slide(prs)
    slide_title(s, "Process flow", "זרימת תהליך", kicker="Components")
    steps = [("icon-mail","black","קליטה"),("icon-cpu","purple","סיווג AI"),
             ("icon-user","black","ניתוב"),("icon-shield-check","green","טיפול"),
             ("icon-chart-bar","black","מדידה")]
    n = len(steps); d = IN(1.3)
    total = IN(11.6); gap = (total - n*d) / (n-1)
    # RTL: first step on the right
    xs = [IN(0.85) + (n-1-i)*(d+gap) for i in range(n)]
    cy = IN(3.5)
    for i,(ic,ink,lab) in enumerate(steps):
        x = xs[i]
        add_oval(s, x, cy, d, ink=ink)
        add_icon(s, ic, x + IN(0.34), cy + IN(0.34), IN(0.62))
        add_text(s, x - IN(0.3), cy + d + IN(0.1), d + IN(0.6), IN(0.5),
                 [[(lab, {"size": 15, "bold": True, "color": INK[ink],
                          "font": F_BODY_HE})]], align=PP_ALIGN.CENTER, rtl=True)
        if i < n-1:
            ax = xs[i] - gap - IN(0.02)
            arr = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, ax, cy + IN(0.5),
                                     gap + IN(0.04), IN(0.3))
            acol = INK["purple"] if steps[i+1][1]=="purple" or ink=="purple" else TXT3
            arr.fill.solid(); arr.fill.fore_color.rgb = acol
            arr.line.fill.background(); arr.shadow.inherit = False
    add_text(s, IN(0.85), IN(5.6), IN(11.6), IN(0.6),
             [[("One AI step (purple) in an otherwise human flow — right-to-left "
                "reading order.", {"size": 14, "color": TXT2})]],
             align=PP_ALIGN.CENTER)

    # ---- 10 · CHARTS (native, editable) ------------------------------
    s = new_slide(prs)
    slide_title(s, "Charts", "תרשימים", kicker="Components · native & editable")
    # donut
    cd = CategoryChartData()
    cd.categories = ["Phone", "Email", "Digital"]
    cd.add_series("Channel", (46, 33, 21))
    gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, IN(0.6), IN(2.1), IN(3.9), IN(4.4), cd)
    ch = gf.chart; ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    pts = ch.plots[0].series[0].points
    for pt, col in zip(pts, [INK["red"], INK["orange"], INK["blue"]]):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = col
    ch.plots[0].has_data_labels = True
    add_text(s, IN(0.6), IN(1.78), IN(3.9), IN(0.3),
             [[("By channel — לפי ערוץ", {"size": 12, "color": TXT2})]])
    # bar
    cd2 = CategoryChartData()
    cd2.categories = ["Sanit.", "Billing", "Eng.", "Welfare"]
    cd2.add_series("Volume", (38, 29, 21, 14))
    gf2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, IN(4.7), IN(2.1), IN(3.9), IN(4.4), cd2)
    ch2 = gf2.chart; ch2.has_legend = False
    b_pts = ch2.plots[0].series[0].points
    for pt, col in zip(b_pts, [INK["red"], INK["orange"], INK["blue"], INK["green"]]):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = col
    add_text(s, IN(4.7), IN(1.78), IN(3.9), IN(0.3),
             [[("By department — לפי מחלקה", {"size": 12, "color": TXT2})]])
    # line
    cd3 = CategoryChartData()
    cd3.categories = ["W1","W2","W3","W4","W5","W6","W7","W8"]
    cd3.add_series("Days", (9, 8.6, 8.1, 7.2, 6.4, 5.5, 4.8, 4.1))
    gf3 = s.shapes.add_chart(XL_CHART_TYPE.LINE, IN(8.8), IN(2.1), IN(3.95), IN(4.4), cd3)
    ch3 = gf3.chart; ch3.has_legend = False
    ser = ch3.plots[0].series[0]
    ser.format.line.color.rgb = INK["green"]; ser.format.line.width = Pt(2.5)
    add_text(s, IN(8.8), IN(1.78), IN(3.9), IN(0.3),
             [[("Response time trend — מגמת זמן מענה", {"size": 12, "color": TXT2})]])

    # ---- 11 · COMPARISON TABLE ---------------------------------------
    s = new_slide(prs)
    slide_title(s, "Comparison table", "טבלת השוואה", kicker="Components")
    data = [
        ("Metric / מדד", "Today / היום", "Target / יעד", "Δ"),
        ("Response time · זמן מענה", "9 days", "<3 days", "-67%"),
        ("Abandoned calls · שיחות ננטשות", "31%", "<10%", "-21pt"),
        ("First-contact resolution · פתרון בפנייה", "22%", "55%", "×2.5"),
        ("AI auto-classification · סיווג AI", "—", "85%", "new"),
        ("Satisfaction · שביעות רצון", "62%", "80%", "+18pt"),
    ]
    rows, colsn = len(data), 4
    tbl_shape = s.shapes.add_table(rows, colsn, IN(0.6), IN(2.0), IN(12.13), IN(4.6))
    tbl = tbl_shape.table
    tbl.columns[0].width = IN(5.2)
    for c in range(1,4): tbl.columns[c].width = IN(2.31)
    val_cols = {1: INK["red"], 2: INK["green"], 3: TXT}
    for r in range(rows):
        tbl.rows[r].height = IN(0.72)
        for c in range(colsn):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER_SUNKEN if r == 0 else PAPER_RAISED
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Pt(8)
            cell.margin_top = cell.margin_bottom = Pt(2)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            is_he = any('֐' <= ch2c <= '׿' for ch2c in data[r][c])
            p.alignment = PP_ALIGN.RIGHT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run(); run.text = data[r][c]
            run.font.size = Pt(13 if r else 13)
            run.font.bold = (r == 0)
            if r == 0:
                run.font.color.rgb = TXT2
            elif c == 0:
                run.font.color.rgb = TXT
            else:
                run.font.color.rgb = val_cols[c]
                run.font.bold = True
            fnt = F_BODY_HE if is_he else F_BODY_LAT
            run.font.name = fnt; _set_cs(run, fnt)

    # ---- 12 · ROADMAP ------------------------------------------------
    s = new_slide(prs)
    slide_title(s, "Roadmap", "מפת דרכים", kicker="Components")
    quarters = [
        ("Q1", "green",  "Done · הושלם",    ["מיפוי תהליכים","בחירת פלטפורמה","פיילוט ראשון"]),
        ("Q2", "orange", "In progress · בתהליך",["הגירת נתונים","חיבור טלפוניה","הדרכת נציגים"]),
        ("Q3", "blue",   "Next · בהמשך",    ["עלייה לאוויר","הפעלת AI","דשבורד הנהלה"]),
        ("Q4", "blue",   "Next · בהמשך",    ["אפליקציית תושבים","טיוטות AI","מדידת יעדים"]),
    ]
    cw = IN(2.95)
    for i,(q,ink,status,items) in enumerate(quarters):
        x = IN(0.6) + i*IN(3.06)
        add_card(s, x, IN(2.1), cw, IN(4.6), ink=ink)
        add_text(s, x + IN(0.25), IN(2.3), cw-IN(0.5), IN(0.6),
                 [[(q, {"size": 30, "bold": True, "color": INK[ink], "font": F_DISP_LAT})]])
        add_text(s, x + IN(0.25), IN(3.0), cw-IN(0.5), IN(0.4),
                 [[(status, {"size": 12, "bold": True, "color": TXT2})]])
        body = [[("• "+it, {"size": 14, "color": TXT, "font": F_BODY_HE})] for it in items]
        add_text(s, x + IN(0.25), IN(3.6), cw-IN(0.5), IN(2.8), body,
                 align=PP_ALIGN.RIGHT, rtl=True, space_after=8)

    # ---- 13 · RISKS & MITIGATION -------------------------------------
    s = new_slide(prs)
    slide_title(s, "Risks & mitigation", "סיכונים ומענים", kicker="Components")
    risks = [
        ("איכות דאטה", ["כפילויות בין קבצים","שדות חסרים"]),
        ("התנגדות עובדים", ["חשש מהחלפת תפקיד","הרגלי Excel"]),
        ("דיוק סיווג AI", ["שפה מעורבת","קטגוריות נדירות"]),
    ]
    cw = IN(3.95)
    for i,(title,items) in enumerate(risks):
        x = IN(0.6) + i*IN(4.13)
        add_card(s, x, IN(2.05), cw, IN(2.7), ink="orange")
        add_icon(s, "icon-warning", x + IN(0.22), IN(2.25), IN(0.5))
        add_text(s, x + IN(0.2), IN(2.9), cw-IN(0.4), IN(0.4),
                 [[(title, {"size": 16, "bold": True, "color": INK["orange"],
                            "font": F_BODY_HE})]], align=PP_ALIGN.RIGHT, rtl=True)
        body = [[("• "+it, {"size": 13, "color": TXT, "font": F_BODY_HE})] for it in items]
        add_text(s, x + IN(0.2), IN(3.4), cw-IN(0.4), IN(1.2), body,
                 align=PP_ALIGN.RIGHT, rtl=True, space_after=6)
    add_card(s, IN(0.6), IN(5.0), IN(12.13), IN(1.7), ink="green")
    label(s, IN(0.85), IN(5.2), IN(6), "Mitigation", "מנגנוני מענה")
    mit = [[("✓ ניקוי דאטה לפני הגירה   ✓ שגרירי שינוי בכל מחלקה   "
             "✓ אדם מאשר כל מענה AI",
             {"size": 15, "color": TXT, "font": F_BODY_HE})]]
    add_text(s, IN(0.85), IN(5.75), IN(11.6), IN(0.8), mit, align=PP_ALIGN.RIGHT,
             rtl=True)

    # ---- 14 · 2x2 MATRIX ---------------------------------------------
    s = new_slide(prs)
    slide_title(s, "2×2 matrix", "מטריצה 2×2", kicker="Components")
    mx, my, msz = IN(3.9), IN(2.1), IN(4.4)
    quad_fill = [PAPER_SUNKEN, PAPER_RAISED, PAPER_RAISED, PAPER_SUNKEN]
    half = msz/2
    coords = [(mx, my),(mx+half, my),(mx, my+half),(mx+half, my+half)]
    for (qx,qy),fl in zip(coords, quad_fill):
        add_rect(s, qx, qy, half, half, fl, line=TXT2, line_w=1.0)
    # sample dots + quadrant labels
    q_labels = [("Quick wins","זכיות מהירות","green",mx+IN(0.2),my+IN(0.2)),
                ("Major bets","הימורים","blue",mx+half+IN(0.2),my+IN(0.2)),
                ("Fill-ins","משני",  "orange",mx+IN(0.2),my+half+IN(0.2)),
                ("Reconsider","לבחון","red",mx+half+IN(0.2),my+half+IN(0.2))]
    for en_l,he_l,ink,lx,ly in q_labels:
        add_text(s, lx, ly, IN(1.9), IN(0.6),
                 [[(en_l, {"size": 13, "bold": True, "color": INK[ink]})],
                  [(he_l, {"size": 13, "color": TXT2, "font": F_BODY_HE})]])
    # axis labels
    add_text(s, mx, my+msz+IN(0.05), msz, IN(0.4),
             [[("Effort → מאמץ", {"size": 12, "bold": True, "color": TXT2})]],
             align=PP_ALIGN.CENTER)
    ylab = add_text(s, mx-IN(1.6), my+half-IN(0.2), IN(1.5), IN(0.4),
             [[("Impact → השפעה", {"size": 12, "bold": True, "color": TXT2})]],
             align=PP_ALIGN.CENTER)
    ylab.rotation = -90
    add_text(s, IN(8.7), IN(2.4), IN(4), IN(4),
             [[("Use for prioritization, risk/impact, or effort/value "
                "trade-offs.", {"size": 15, "color": TXT2})],
              [("", {"size":8,"color":TXT2})],
              [("לתעדוף, סיכון מול השפעה, או מאמץ מול ערך.",
                {"size": 15, "color": TXT2, "font": F_BODY_HE})]],
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    for p in ylab.text_frame.paragraphs: pass

    # ---- 15 · CLOSING ------------------------------------------------
    s = new_slide(prs)
    add_rect(s, 0, IN(7.22), EMU_W, IN(0.28), ACCENT)
    add_icon(s, "icon-sparkle", IN(0.7), IN(0.85), IN(0.9))
    add_text(s, IN(0.7), IN(2.1), IN(12), IN(1.2),
             [[("Build on-brand.", {"size": 52, "bold": True, "color": TXT,
                                    "font": F_DISP_LAT})]])
    add_text(s, IN(0.7), IN(3.3), IN(12), IN(0.8),
             [[("בנו במותג — עברית ואנגלית, שש דיו, אותה שפה חזותית",
                {"size": 26, "color": ACCENT, "font": F_DISP_HE})]],
             align=PP_ALIGN.RIGHT, rtl=True)
    recap = ("Colors · Typography · Labels · Icons · KPIs · Problem/Solution · "
             "Flow · Charts · Table · Roadmap · Risks · Matrix")
    add_text(s, IN(0.72), IN(4.5), IN(12), IN(1.2),
             [[(recap, {"size": 15, "color": TXT2})]])
    add_text(s, IN(0.72), IN(5.6), IN(12), IN(0.5),
             [[("Workshop Design System · github.com/noamh98/workshop-design-system",
                {"size": 13, "color": TXT3, "font": F_MONO})]])
    dx = IN(0.72)
    for k in ["black","blue","red","green","orange","purple"]:
        add_oval(s, dx, IN(6.4), IN(0.3), ink=k, fill=INK[k], line_w=0)
        dx += IN(0.5)

    prs.save(out_path)
    print("saved", out_path, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")

if __name__ == "__main__":
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    out = args[0] if args else "workshop-design-system-template.pptx"
    if "--icons" in sys.argv:
        ICON_DIR = Path(sys.argv[sys.argv.index("--icons")+1])
    build(out)
